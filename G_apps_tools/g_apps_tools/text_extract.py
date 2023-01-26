import collections.abc
from pptx import Presentation
import fitz
import pandas as pd

def extract_from_pptx(doc_name):
    texts = ""
    try:
        presentation = Presentation('./doc/' + doc_name)
    except:
        return "ない"

    for sld in presentation.slides:
        for shape in sld.shapes:
            # shapeに含まれるテキストデータを抽出
            if shape.has_text_frame:
                for text in shape.text.splitlines():
                    texts += text
            # tableに含まれるテキストデータを抽出
            if shape.has_table: 
                for cell in shape.table.iter_cells():
                    for text in cell.text.splitlines():
                        texts += text
    return texts

def extract_from_pdf(doc_name):
    texts = ""
    try:
        pdf_file = fitz.open('./doc/' + doc_name)
    except:
        return "わからん"

    for page in pdf_file:
        text = page.get_text()
        text = text.replace('\n', '')
        texts += text
    return texts

def check_filetype(doc_name):
    if doc_name.endswith('.pptx'):
        doc_text = extract_from_pptx(doc_name)
    elif doc_name.endswith('.pdf'):
        doc_text = extract_from_pdf(doc_name)
    else:
        doc_text = "none"
    return doc_text


def txt_ext(doc_list):
    doc_list["doc_text"] = doc_list["doc_title"].apply(check_filetype)
    return doc_list

data = pd.read_csv("./data/update_doc_data.csv")
update_data = txt_ext(data)
update_data.to_csv("./data/update_doc_datav2.csv", index = False)